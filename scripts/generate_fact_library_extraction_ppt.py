from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PPT = Path(
    "/Users/caixudong/Downloads/zhilian-robot/docs/reports/fact_library_extraction_concept_layer_presentation.pptx"
)


# Visual system
BG = RGBColor(244, 247, 252)
NAVY = RGBColor(20, 36, 66)
BLUE = RGBColor(39, 107, 190)
MINT = RGBColor(94, 173, 164)
TEXT = RGBColor(45, 55, 72)
MUTED = RGBColor(89, 102, 124)
CARD = RGBColor(255, 255, 255)


def _font(run, size=16, bold=False, color=TEXT):
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_box_text(shape, text, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _font(r, size=size, bold=bold, color=color)


def _add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _add_header_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    _font(r, size=20, bold=True, color=RGBColor(255, 255, 255))


def _add_title(slide, title, subtitle=""):
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.25), Inches(11.8), Inches(1.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    _font(r, size=30, bold=True, color=NAVY)

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.78), Inches(2.2), Inches(10.0), Inches(0.9))
        stf = sb.text_frame
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        _font(sr, size=16, color=MUTED)


def _add_bullets(slide, bullets, left=0.9, top=1.8, width=11.6, height=4.8, size=19):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, txt in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.bullet = True
        r = p.add_run()
        r.text = txt
        _font(r, size=size, color=TEXT)


def _add_card(slide, left, top, width, height, title, lines, tone=CARD):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = tone
    card.line.color.rgb = RGBColor(214, 222, 235)

    title_box = slide.shapes.add_textbox(Inches(left + 0.22), Inches(top + 0.16), Inches(width - 0.4), Inches(0.35))
    _set_box_text(title_box, title, size=16, bold=True, color=NAVY)

    body = slide.shapes.add_textbox(Inches(left + 0.24), Inches(top + 0.58), Inches(width - 0.42), Inches(height - 0.72))
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.bullet = True
        r = p.add_run()
        r.text = line
        _font(r, size=13, color=TEXT)


def _cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)

    hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(230, 238, 250)
    hero.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.0))
    accent.fill.solid()
    accent.fill.fore_color.rgb = NAVY
    accent.line.fill.background()

    _add_title(
        slide,
        "试点知识库抽取工作与概念层引入方案",
        "基于 Fact Library Pipeline + OpenSPG/KAG | 2026-03-17",
    )

    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(3.35), Inches(5.6), Inches(0.6))
    tag.fill.solid()
    tag.fill.fore_color.rgb = RGBColor(255, 255, 255)
    tag.line.color.rgb = RGBColor(187, 203, 227)
    _set_box_text(tag, "目标：讲清当前抽取流程 + 概念层落地路径", size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)


def _agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_header_bar(slide, "汇报目录")
    _add_card(slide, 0.9, 1.45, 5.7, 4.9, "第一部分：当前抽取工作", [
        "当前抽取目标与边界",
        "实体 / 文本 / 关系抽取流程",
        "规则设计与质量控制",
        "当前产出能力与短板",
    ])
    _add_card(slide, 6.75, 1.45, 5.7, 4.9, "第二部分：概念层引入", [
        "概念层的定义与价值",
        "面向现有知识库的引入方案",
        "分阶段改造路线",
        "预期收益与下阶段动作",
    ])


def _scope_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_header_bar(slide, "一、当前抽取工作")
    _add_bullets(slide, [
        "当前目标：先稳定构建“基础事实知识层”，保证数据可持续入库。",
        "当前边界：以结构化抽取为主，不在首阶段做大规模 LLM 语义抽取。",
        "当前输出：`entities/` + `support/` + `texts/` + `relations/` 四类标准化文件。",
        "当前定位：先把“事实装进去”，再逐步增强“概念归纳与事件推理”。",
    ], top=1.55, size=20)


def _pipeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_header_bar(slide, "抽取流程（端到端）")

    steps = [
        "读取原始表",
        "规则筛选",
        "ID去重",
        "字段分层",
        "文本拼接",
        "关系物化",
        "图谱导入",
    ]
    left = 0.65
    top = 2.2
    w = 1.65
    h = 1.2
    gap = 0.18
    for i, s in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD
        box.line.color.rgb = RGBColor(194, 209, 232)
        _set_box_text(box, s, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(left + w + 0.02),
                Inches(top + 0.38),
                Inches(gap - 0.04),
                Inches(0.45),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BLUE
            arrow.line.fill.background()
        left += w + gap

    note = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.8), Inches(1.9))
    tf = note.text_frame
    tf.clear()
    lines = [
        "关键点1：`keep_columns` 保内部工作字段，`export_columns` 保最终实体字段。",
        "关键点2：关系只连到本轮保留实体，避免悬空边。",
        "关键点3：文本材料与实体导出解耦，后续可单独接 KAG/OpenKS 增强抽取。",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.bullet = True
        r = p.add_run()
        r.text = line
        _font(r, size=15, color=TEXT)


def _rules_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_header_bar(slide, "当前规则设计（示例）")

    _add_card(slide, 0.75, 1.5, 3.0, 4.9, "实体筛选规则", [
        "Company：状态正常",
        "Institution：domain非空",
        "Project：2018年后",
        "Patent：授权日期非空",
        "Article：2018年后",
    ])
    _add_card(slide, 3.95, 1.5, 3.0, 4.9, "字段分层规则", [
        "export_columns：主实体输出",
        "keep_columns：关系/文本工作字段",
        "text_fields：拼接增强材料",
        "按 id 去重，空 id 丢弃",
    ])
    _add_card(slide, 7.15, 1.5, 5.4, 4.9, "关系抽取规则", [
        "建立实体名称索引",
        "多值拆分 + 名称标准化",
        "企业/机构启发式去歧义",
        "匹配到实体ID后物化关系",
        "统一输出 s_id/s_type/p/o_id/o_type/properties",
        "导出 relation_summary 与 unmatched 列表",
    ])


def _result_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_header_bar(slide, "当前抽取结果与评价")
    _add_bullets(slide, [
        "当前已稳定产出：主实体、辅助表、文本材料、显式关系四类结果。",
        "关系覆盖：项目承担、专利归属、成果归属、文献作者、榜单包含等核心事实链。",
        "优势：规则清晰、可解释、速度快、适配大批量结构化数据。",
        "不足：精确匹配依赖较强，概念层与事件层尚未系统化建设。",
    ], top=1.55, size=19)

    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.85), Inches(11.5), Inches(1.0))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(226, 239, 255)
    band.line.color.rgb = RGBColor(181, 205, 236)
    _set_box_text(band, "结论：基础事实层已跑通，下一步应重点补“概念层”以提升组织能力与推理能力。", size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def _concept_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_header_bar(slide, "二、概念层引入")

    left = 0.9
    top = 1.65
    width = 3.9
    height = 3.7
    _add_card(slide, left, top, width, height, "实体层（Fact）", [
        "存具体对象",
        "企业、人物、项目、专利",
        "回答“有什么”",
    ], tone=RGBColor(250, 252, 255))
    _add_card(slide, left + 4.2, top, width, height, "事件层（Event）", [
        "存具体发生的事",
        "授权、立项、发布、中标",
        "回答“发生了什么”",
    ], tone=RGBColor(248, 255, 253))
    _add_card(slide, left + 8.4, top, width, height, "概念层（Concept）", [
        "存分类体系与层级",
        "行业、IPC、学科、标准状态",
        "回答“属于哪一类”",
    ], tone=RGBColor(255, 250, 244))

    hint = slide.shapes.add_textbox(Inches(0.95), Inches(5.65), Inches(11.6), Inches(0.9))
    tf = hint.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "概念层不是普通标签，而是可复用、可传播、可推理的分类语义层。"
    _font(r, size=17, bold=True, color=NAVY)


def _integration_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_header_bar(slide, "概念层如何引入当前知识库")

    _add_card(slide, 0.8, 1.5, 3.9, 4.8, "阶段1：先补静态概念", [
        "CompanyTaxonomy",
        "IndustryTaxonomy",
        "IPCTaxonomy",
        "SubjectTaxonomy",
        "StandardStatusTaxonomy",
    ])
    _add_card(slide, 4.95, 1.5, 3.9, 4.8, "阶段2：建立实体-概念挂载", [
        "Company -> belongTo -> CompanyTaxonomy",
        "Patent -> belongTo -> IPCTaxonomy",
        "Article -> belongTo -> SubjectTaxonomy",
        "Standard -> belongTo -> StandardStatusTaxonomy",
    ])
    _add_card(slide, 9.1, 1.5, 3.4, 4.8, "阶段3：升级事件与规则", [
        "补事件类型",
        "接入概念推导",
        "增强检索与问答",
        "逐步引入推理规则",
    ])


def _closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_header_bar(slide, "结论与下一步")
    _add_bullets(slide, [
        "当前阶段：基础事实知识抽取流程已跑通，可稳定导入 OpenSPG/KAG。",
        "核心短板：缺少概念层，知识仍偏“事实仓库”。",
        "执行建议：先引入静态概念层，再补事件层和规则层，分阶段演进。",
        "预期收益：提高检索精度、抽取边界清晰度、Schema 可维护性和后续推理能力。",
    ], top=1.55, size=20)

    end = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.0))
    end.fill.solid()
    end.fill.fore_color.rgb = RGBColor(235, 244, 255)
    end.line.color.rgb = RGBColor(180, 204, 233)
    _set_box_text(end, "建议先用 1 个迭代完成“概念层最小闭环”验证。", size=17, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _cover_slide(prs)
    _agenda_slide(prs)
    _scope_slide(prs)
    _pipeline_slide(prs)
    _rules_slide(prs)
    _result_slide(prs)
    _concept_intro_slide(prs)
    _integration_slide(prs)
    _closing_slide(prs)

    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPT)
    print(OUT_PPT)
    print("slides=9")


if __name__ == "__main__":
    build()
