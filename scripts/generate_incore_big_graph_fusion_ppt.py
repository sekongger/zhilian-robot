from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PPT = Path(
    "/Users/caixudong/Downloads/zhilian-robot/docs/reports/incore_big_graph_fusion_layer_presentation.pptx"
)


BG = RGBColor(245, 247, 251)
NAVY = RGBColor(18, 34, 64)
BLUE = RGBColor(40, 102, 184)
CYAN = RGBColor(83, 165, 190)
TEAL = RGBColor(86, 168, 152)
TEXT = RGBColor(47, 60, 79)
MUTED = RGBColor(103, 117, 138)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(208, 217, 231)
PALE_BLUE = RGBColor(232, 240, 252)
PALE_GREEN = RGBColor(235, 246, 242)
PALE_CYAN = RGBColor(235, 246, 250)
PALE_GOLD = RGBColor(250, 245, 231)
WHITE = RGBColor(255, 255, 255)


def _font(run, size=16, bold=False, color=TEXT):
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _set_text(
    shape,
    text,
    size=16,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _font(r, size=size, bold=bold, color=color)


def _add_paragraphs(shape, lines, size=15, color=TEXT, bullet=True):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.bullet = bullet
        r = p.add_run()
        r.text = line
        _font(r, size=size, color=color)


def _header(slide, title, kicker=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.16), Inches(10.0), Inches(0.45))
    _set_text(tb, title, size=22, bold=True, color=WHITE)

    if kicker:
        kb = slide.shapes.add_textbox(Inches(10.5), Inches(0.18), Inches(2.2), Inches(0.35))
        _set_text(kb, kicker, size=11, color=WHITE, align=PP_ALIGN.RIGHT)


def _title_block(slide, title, subtitle):
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.1), Inches(11.4), Inches(0.8))
    _set_text(tb, title, size=29, bold=True, color=NAVY)
    sb = slide.shapes.add_textbox(Inches(0.78), Inches(1.95), Inches(11.2), Inches(0.7))
    _set_text(sb, subtitle, size=15, color=MUTED)


def _tag(slide, left, top, width, text, fill):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.5)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    _set_text(shp, text, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def _card(slide, left, top, width, height, title, body_lines, fill=CARD, title_color=NAVY):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = LINE

    tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.14), Inches(width - 0.35), Inches(0.35))
    _set_text(tb, title, size=16, bold=True, color=title_color)

    body = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.58), Inches(width - 0.35), Inches(height - 0.72)
    )
    _add_paragraphs(body, body_lines, size=13, color=TEXT, bullet=True)


def _number_badge(slide, left, top, text, fill=BLUE):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.42), Inches(0.42))
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.fill.background()
    _set_text(c, text, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def _connector(slide, x1, y1, x2, y2, color=BLUE, weight=2.25):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(weight)


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)

    hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(230, 238, 250)
    hero.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.0))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = NAVY
    top_bar.line.fill.background()

    _title_block(
        slide,
        "IncCore 大图融合层图谱构建方案",
        "基于统一 schema 的常识层 + 事实层融合建图，支撑后续推理、抽取与产业网链计算 | 2026-03-22",
    )

    _tag(slide, 0.8, 3.0, 2.4, "统一语义骨架", PALE_BLUE)
    _tag(slide, 3.4, 3.0, 2.4, "多源融合入图", PALE_GREEN)
    _tag(slide, 6.0, 3.0, 2.4, "概念 + 事件协同", PALE_CYAN)
    _tag(slide, 8.6, 3.0, 2.4, "推理友好结构", PALE_GOLD)

    summary = slide.shapes.add_textbox(Inches(0.85), Inches(4.0), Inches(11.2), Inches(2.0))
    _add_paragraphs(
        summary,
        [
            "目标不是新增一张孤立图谱，而是用 IncCore.schema 统一承载常识层与资讯/研报事实层。",
            "核心抓手是三件事：主实体对齐、关系与属性融合、概念层与事件层协同建模。",
            "最终产物是一张可追溯、可推理、可继续被 KAG/OpenKS 利用的统一 OpenSPG 大图。",
        ],
        size=18,
        color=TEXT,
        bullet=True,
    )


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _header(slide, "方案总览", "结论先行")

    _card(
        slide,
        0.7,
        1.4,
        3.85,
        4.9,
        "为什么要做融合层",
        [
            "常识层解决“主体是谁、长期关系是什么”。",
            "事实层解决“最近发生了什么、如何变化”。",
            "统一大图解决“如何把背景、事件和证据放进同一世界模型”。",
        ],
        fill=PALE_BLUE,
    )
    _card(
        slide,
        4.75,
        1.4,
        3.85,
        4.9,
        "方案主线",
        [
            "以 IncCore.schema 作为唯一语义边界。",
            "先做统一接入、主实体对齐，再做关系融合与冲突消解。",
            "在统一大图中引入概念层和事件层，服务后续计算。",
        ],
        fill=PALE_GREEN,
    )
    _card(
        slide,
        8.8,
        1.4,
        3.85,
        4.9,
        "最终结果",
        [
            "一张可统一管理常识、事实、证据的 OpenSPG 大图。",
            "一套可支持产业网链推理、问答、风险预警的计算底座。",
            "一条可持续接入资讯、研报等新来源的融合 pipeline。",
        ],
        fill=PALE_CYAN,
    )

    bottom = slide.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(11.8), Inches(0.5))
    _set_text(bottom, "一句话：统一 schema 是边界，主实体融合是入口，概念层和事件层是推理基础。", size=18, bold=True, color=NAVY)


def slide_schema(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _header(slide, "IncCore.schema 在融合层中的作用")

    _card(
        slide,
        0.7,
        1.45,
        3.9,
        4.7,
        "概念层骨架",
        [
            "IndustrySector、CompanyCategory、ProductCategory、TechnologyCategory。",
            "PersonCategory、OrganizationCategory、RegionCategory、EventCategory、TermCategory。",
            "这层不是标签，而是实例归类、概念传播和抽取约束的承接层。",
        ],
        fill=PALE_BLUE,
    )
    _card(
        slide,
        4.75,
        1.45,
        3.9,
        4.7,
        "实体层骨架",
        [
            "IndustryActor 统一承载 Company、Organization、Person。",
            "Technology、ProductObject、Region 提供产业网链中的关键实体对象。",
            "Document、Chunk、DataSource 为事实证据和可追溯能力提供承载对象。",
        ],
        fill=PALE_GREEN,
    )
    _card(
        slide,
        8.8,
        1.45,
        3.9,
        4.7,
        "事件层骨架",
        [
            "已有 Event、GovernmentPublishPolicyEvent、CompanyCooperationEvent、CompanyFinancingEvent。",
            "说明 schema 已经允许把“动态事实”接入统一大图，而不是只存静态实体关系。",
            "后续只需要在此基础上继续扩展事件族即可。",
        ],
        fill=PALE_CYAN,
    )

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.55))
    _set_text(foot, "一句话：IncCore.schema 已经具备“概念 + 实体 + 事件 + 证据”的统一建模起点。", size=18, bold=True, color=NAVY)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _header(slide, "大图融合层总体架构")

    nodes = [
        (0.7, 2.15, 2.0, 1.05, "常识层输入", PALE_BLUE),
        (0.7, 3.65, 2.0, 1.05, "事实层输入", PALE_GREEN),
        (0.7, 5.15, 2.0, 1.05, "词表/规则输入", PALE_GOLD),
        (3.2, 3.0, 2.2, 1.1, "统一接入与标准化", CARD),
        (5.7, 3.0, 2.15, 1.1, "主实体对齐\n与主键归一", CARD),
        (8.05, 3.0, 2.15, 1.1, "关系融合\n与冲突消解", CARD),
        (10.4, 3.0, 2.1, 1.1, "概念挂载\n与事件建模", CARD),
        (10.45, 5.0, 2.05, 0.95, "OpenSPG\n统一大图", PALE_CYAN),
    ]
    for left, top, width, height, text, fill in nodes:
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = LINE
        _set_text(shp, text, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    _connector(slide, 2.7, 2.68, 3.18, 3.55)
    _connector(slide, 2.7, 4.18, 3.18, 3.55)
    _connector(slide, 2.7, 5.68, 3.18, 3.55)
    _connector(slide, 5.4, 3.55, 5.68, 3.55)
    _connector(slide, 7.85, 3.55, 8.03, 3.55)
    _connector(slide, 10.2, 3.55, 10.38, 3.55)
    _connector(slide, 11.45, 4.12, 11.45, 4.98)

    note = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.6))
    _set_text(note, "一句话：融合层不是“直接灌图”，而是先统一语义、再对齐主实体、再落概念和事件。", size=18, bold=True, color=NAVY)


def slide_fusion_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _header(slide, "多源融合流程")

    steps = [
        ("1", "统一 DTO", ["所有源数据先转成统一融合对象，不能直接按原字段入图。"]),
        ("2", "主键归一", ["企业以信用代码优先，人物结合机构与职务，事件采用软主键。"]),
        ("3", "主实体对齐", ["规则对齐 -> 名称归一 -> 上下文辅助消歧，分层完成融合。"]),
        ("4", "属性与关系融合", ["静态属性按权威源优先，描述型信息保留多来源，关系保留证据和置信度。"]),
        ("5", "概念挂载与事件化", ["分类信息上提到概念层，动态变化优先进入事件层。"]),
    ]
    top = 1.55
    for idx, (num, title, lines) in enumerate(steps):
        y = top + idx * 0.95
        _number_badge(slide, 0.75, y + 0.08, num, fill=BLUE if idx < 3 else TEAL)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.28), Inches(y), Inches(11.0), Inches(0.72)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = LINE
        title_box = slide.shapes.add_textbox(Inches(1.55), Inches(y + 0.12), Inches(2.2), Inches(0.25))
        _set_text(title_box, title, size=15, bold=True, color=NAVY)
        desc_box = slide.shapes.add_textbox(Inches(3.55), Inches(y + 0.1), Inches(8.4), Inches(0.4))
        _add_paragraphs(desc_box, lines, size=13, color=TEXT, bullet=False)

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.45))
    _set_text(foot, "一句话：先统一对象，再融合主实体，最后把静态归常识、动态归事件。", size=18, bold=True, color=NAVY)


def slide_alignment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _header(slide, "主实体对齐与冲突消解")

    _card(
        slide,
        0.7,
        1.45,
        3.8,
        4.95,
        "主实体对齐",
        [
            "Company：信用代码优先，名称归一为辅。",
            "Organization：标准机构名 + 区域，多源编码作为补充。",
            "Person：姓名 + 机构 + 职务，必要时结合研究方向。",
            "Product / Technology：标准名优先，品牌/型号/术语作为补充锚点。",
        ],
        fill=PALE_BLUE,
    )
    _card(
        slide,
        4.75,
        1.45,
        3.8,
        4.95,
        "冲突消解",
        [
            "身份型字段：权威源优先，冲突值保留候选与来源。",
            "描述型字段：不做覆盖，保留多来源摘要与时间。",
            "关系型字段：以边为中心管理，保留来源、时间和置信度。",
            "时态型字段：尽量事件化，不直接覆盖静态属性。",
        ],
        fill=PALE_GREEN,
    )
    _card(
        slide,
        8.8,
        1.45,
        3.8,
        4.95,
        "来源优先级",
        [
            "1. 政务/监管来源",
            "2. 企业官网与官方公告",
            "3. 权威数据库",
            "4. 主流财经与行业媒体",
            "5. 一般资讯与转载来源",
        ],
        fill=PALE_GOLD,
    )

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.52), Inches(11.7), Inches(0.48))
    _set_text(foot, "一句话：融合层的关键不是“全并进来”，而是先找到统一主节点，再有控制地保留差异。", size=18, bold=True, color=NAVY)


def slide_concepts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _header(slide, "概念层设计")

    _card(
        slide,
        0.7,
        1.45,
        3.0,
        4.9,
        "概念层做什么",
        [
            "统一不同来源的分类口径。",
            "承接实例归类和抽象传播。",
            "作为推理约束和抽取边界。",
            "为产业网链推理提供概念骨架。",
        ],
        fill=PALE_BLUE,
    )
    _card(
        slide,
        3.95,
        1.45,
        4.2,
        4.9,
        "当前 schema 已有概念",
        [
            "IndustrySector",
            "CompanyCategory / OrganizationCategory / PersonCategory",
            "ProductCategory / TechnologyCategory",
            "RegionCategory / EventCategory / TermCategory",
            "已足够支撑第一轮概念挂载。",
        ],
        fill=PALE_GREEN,
    )
    _card(
        slide,
        8.4,
        1.45,
        4.2,
        4.9,
        "概念如何获取",
        [
            "结构化字段直接映射。",
            "标准词表与人工维护目录导入。",
            "从资讯、研报文本中做概念抽取。",
            "基于图结构关系做归纳补充。",
        ],
        fill=PALE_CYAN,
    )

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.48))
    _set_text(foot, "一句话：概念层不是标签层，而是后续检索、抽取、传播和解释的抽象层。", size=18, bold=True, color=NAVY)


def slide_events(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _header(slide, "事件层设计")

    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(2.35), Inches(1.0)
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = PALE_GREEN
    box1.line.color.rgb = LINE
    _set_text(box1, "事实来源\n资讯 / 研报 / 公告", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(1.55), Inches(3.0), Inches(1.35)
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = CARD
    box2.line.color.rgb = LINE
    _set_text(box2, "统一事件对象\nEvent / PolicyEvent / CooperationEvent / FinancingEvent", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    box3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.9), Inches(1.55), Inches(2.7), Inches(1.35)
    )
    box3.fill.solid()
    box3.fill.fore_color.rgb = PALE_CYAN
    box3.line.color.rgb = LINE
    _set_text(box3, "事件聚合\n按“类型 + 主体 + 客体 + 时间 + 地点”去重", size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    box4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.95), Inches(1.7), Inches(1.6), Inches(1.0)
    )
    box4.fill.solid()
    box4.fill.fore_color.rgb = PALE_GOLD
    box4.line.color.rgb = LINE
    _set_text(box4, "证据回连\nDocument / Chunk", size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    _connector(slide, 3.16, 2.2, 3.95, 2.2)
    _connector(slide, 7.03, 2.2, 7.88, 2.2)
    _connector(slide, 10.62, 2.2, 10.93, 2.2)

    _card(
        slide,
        0.8,
        3.45,
        5.7,
        2.45,
        "当前 schema 已能承载的事件",
        [
            "GovernmentPublishPolicyEvent：政策发布、监管措施、区域政策。",
            "CompanyCooperationEvent：签约合作、联合研发、战略协作。",
            "CompanyFinancingEvent：融资、投资、轮次和金额变化。",
            "这三类覆盖了资讯与研报中最常见、最有价值的一批产业事件。",
        ],
        fill=PALE_BLUE,
    )
    _card(
        slide,
        6.75,
        3.45,
        5.8,
        2.45,
        "下一轮建议扩展",
        [
            "投资建设事件、产品发布事件、技术突破事件、产能扩张事件、风险异常事件。",
            "统一补充事件摘要、事件时间、影响范围、相关产品/技术、置信度等公共属性。",
            "这样事件层才能直接服务产业链传导分析和问答解释。",
        ],
        fill=PALE_GREEN,
    )

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.42), Inches(11.7), Inches(0.52))
    _set_text(foot, "一句话：事件层承载动态事实，证据层负责可追溯，二者一起解决“最近发生了什么”。", size=18, bold=True, color=NAVY)


def slide_reasoning(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _header(slide, "如何服务后续推理与抽取")

    _card(
        slide,
        0.7,
        1.45,
        5.85,
        4.95,
        "服务推理",
        [
            "概念层先缩小候选范围，再落到实例层做产业链检索。",
            "常识层提供长期背景关系，事件层提供动态变化输入。",
            "可以支持上下游传导、区域影响、政策影响、融资链路和风险预警。",
            "证据层让每条推理结果都能回到原文和原片段。",
        ],
        fill=PALE_BLUE,
    )
    _card(
        slide,
        6.75,
        1.45,
        5.85,
        4.95,
        "服务抽取增强",
        [
            "概念层可反向作为 KAG/OpenKS 的抽取目标边界。",
            "术语、产品、技术、事件分类都可以变成抽取词表和 schema 约束。",
            "统一大图还能为实体消歧、概念归类和事件聚合提供背景知识。",
            "后续抽取不是孤立处理文本，而是有统一语义世界做对齐。",
        ],
        fill=PALE_GREEN,
    )

    bottom = slide.shapes.add_textbox(Inches(0.8), Inches(6.52), Inches(11.7), Inches(0.5))
    _set_text(bottom, "一句话：统一大图既是计算底座，也是后续抽取系统的背景知识底座。", size=18, bold=True, color=NAVY)


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _header(slide, "实施路径与近期动作")

    phases = [
        ("Phase 1", "统一骨架落图", ["先把常识层核心实体、概念和最小事件集统一入图。"], PALE_BLUE),
        ("Phase 2", "主实体融合", ["打通企业、机构、人物、产品、技术的跨源对齐和证据回连。"], PALE_GREEN),
        ("Phase 3", "概念层增强", ["建立产业、产品、技术、事件分类体系，形成概念挂载闭环。"], PALE_CYAN),
        ("Phase 4", "推理友好化", ["扩展事件类型和概念关系，支撑传导分析、问答与风险预警。"], PALE_GOLD),
    ]
    left = 0.7
    top = 2.0
    width = 2.9
    for idx, (phase, title, lines, fill) in enumerate(phases):
        _card(slide, left + idx * 3.05, top, width, 3.25, f"{phase} | {title}", lines, fill=fill)

    action = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.6), Inches(0.85))
    _add_paragraphs(
        action,
        [
            "建议优先输出两份落地文件：一份是“源数据 -> IncCore.schema”的对象映射表，一份是“IncCore.schema v2”的扩展草案。",
            "这样下一步就能直接进入融合 pipeline 和 OpenSPG 导入链的实施阶段。",
        ],
        size=16,
        color=TEXT,
        bullet=True,
    )


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover(prs)
    slide_conclusion(prs)
    slide_schema(prs)
    slide_architecture(prs)
    slide_fusion_flow(prs)
    slide_alignment(prs)
    slide_concepts(prs)
    slide_events(prs)
    slide_reasoning(prs)
    slide_roadmap(prs)

    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(f"written {OUT_PPT}")


if __name__ == "__main__":
    build()
