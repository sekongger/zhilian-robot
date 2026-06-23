from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PPT = Path(
    "/Users/caixudong/Downloads/zhilian-robot/docs/reports/incore_big_graph_overview.pptx"
)


BG = RGBColor(245, 247, 251)
NAVY = RGBColor(16, 34, 63)
BLUE = RGBColor(39, 103, 190)
TEAL = RGBColor(84, 167, 160)
GOLD = RGBColor(208, 158, 54)
TEXT = RGBColor(43, 55, 76)
MUTED = RGBColor(101, 113, 134)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(208, 217, 231)
CARD = RGBColor(255, 255, 255)
PALE_BLUE = RGBColor(234, 241, 251)
PALE_GREEN = RGBColor(235, 246, 242)
PALE_GOLD = RGBColor(251, 245, 232)
PALE_CYAN = RGBColor(235, 247, 249)


def _font(run, size=16, bold=False, color=TEXT):
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_text(shape, text, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _font(r, size=size, bold=bold, color=color)


def _add_paragraphs(shape, lines, size=14, color=TEXT, bullet=True):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.bullet = bullet
        r = p.add_run()
        r.text = line
        _font(r, size=size, color=color)


def _background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _header(slide, title, kicker=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.16), Inches(9.5), Inches(0.45))
    _set_text(title_box, title, size=22, bold=True, color=WHITE)

    if kicker:
        kick_box = slide.shapes.add_textbox(Inches(10.4), Inches(0.18), Inches(2.2), Inches(0.3))
        _set_text(kick_box, kicker, size=11, color=WHITE, align=PP_ALIGN.RIGHT)


def _rounded_box(slide, left, top, width, height, text, fill=CARD, size=15, bold=True, color=NAVY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    _set_text(shape, text, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER)
    return shape


def _section(slide, left, top, width, height, title, fill):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = LINE

    title_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.08), Inches(width - 0.3), Inches(0.3))
    _set_text(title_box, title, size=16, bold=True, color=NAVY)
    return box


def _arrow(slide, x1, y1, x2, y2, color=BLUE, width=2.2):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _background(slide)

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.0))
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.2), Inches(1.0))
    _set_text(title, "IncCore 大图设计总览", size=29, bold=True, color=NAVY, valign=MSO_ANCHOR.TOP)

    subtitle = slide.shapes.add_textbox(Inches(0.82), Inches(2.35), Inches(11.3), Inches(0.8))
    _set_text(
        subtitle,
        "用一张图解释：常识层、事实层、概念层、事件层如何统一进入 OpenSPG 大图，并服务后续推理与抽取。",
        size=16,
        color=MUTED,
        valign=MSO_ANCHOR.TOP,
    )

    _rounded_box(slide, 0.85, 3.3, 2.4, 0.6, "统一 schema", fill=PALE_BLUE, size=14)
    _rounded_box(slide, 3.55, 3.3, 2.4, 0.6, "多源融合", fill=PALE_GREEN, size=14)
    _rounded_box(slide, 6.25, 3.3, 2.4, 0.6, "概念 + 事件协同", fill=PALE_CYAN, size=14)
    _rounded_box(slide, 8.95, 3.3, 2.4, 0.6, "推理友好结构", fill=PALE_GOLD, size=14)

    body = slide.shapes.add_textbox(Inches(0.9), Inches(4.35), Inches(11.2), Inches(1.5))
    _add_paragraphs(
        body,
        [
            "左边是常识层、资讯、研报和标准词表等输入源。",
            "中间是统一接入、主实体对齐、冲突消解、概念挂载和事件建模。",
            "右边是一张包含概念层、常识实体层、事件层、证据层的统一 OpenSPG 大图。",
        ],
        size=18,
        color=TEXT,
        bullet=True,
    )


def overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _background(slide)
    _header(slide, "IncCore 大图整体设计图", "总图")

    # Column containers
    _section(slide, 0.35, 1.1, 2.15, 5.95, "输入数据", PALE_BLUE)
    _section(slide, 2.7, 1.1, 2.15, 5.95, "前置处理", PALE_GREEN)
    _section(slide, 5.05, 1.1, 2.8, 5.95, "大图融合层", PALE_GOLD)
    _section(slide, 8.05, 1.1, 2.45, 5.95, "统一产业知识大图", PALE_CYAN)
    _section(slide, 10.7, 1.1, 2.25, 5.95, "知识计算与应用", PALE_BLUE)

    # Input boxes
    _rounded_box(slide, 0.55, 1.75, 1.75, 1.0, "常识层数据\n企业 / 机构 / 人物\n产品 / 技术 / 区域", fill=WHITE, size=13)
    _rounded_box(slide, 0.55, 3.0, 1.75, 1.0, "事实层数据\n资讯 / 研报 / 公告\n公众号 / 数据接口", fill=WHITE, size=13)
    _rounded_box(slide, 0.55, 4.25, 1.75, 1.0, "标准词表与规则库\n产业 / 产品 / 技术\n事件 / 区域分类", fill=WHITE, size=13)

    # Pre-processing
    _rounded_box(slide, 2.95, 1.75, 1.65, 1.0, "采集标化层\n清洗 / 去重\n标准字段", fill=WHITE, size=13)
    _rounded_box(slide, 2.95, 3.0, 1.65, 1.0, "事件聚合层\nGraphiti 抽取\n归一 / 聚类", fill=WHITE, size=13)
    _rounded_box(slide, 2.95, 4.25, 1.65, 1.0, "常识层加工\n基础事实抽取\n关系物化", fill=WHITE, size=13)

    # Fusion boxes
    _rounded_box(slide, 5.3, 1.55, 2.3, 0.78, "统一语义映射\n全部来源先映射到 IncCore.schema", fill=WHITE, size=13)
    _rounded_box(slide, 5.3, 2.55, 2.3, 0.78, "主实体对齐\n统一主键 / 别名归一 / 消歧", fill=WHITE, size=13)
    _rounded_box(slide, 5.3, 3.55, 2.3, 0.78, "关系融合与冲突消解\n多源关系合并 / 权威源优先", fill=WHITE, size=13)
    _rounded_box(slide, 5.3, 4.55, 2.3, 0.78, "概念挂载 + 事件建模\n分类归一 / 事件化表达 / 证据回连", fill=WHITE, size=13)

    # Graph layers
    _rounded_box(slide, 8.28, 1.6, 1.98, 0.8, "概念层\nIndustry / Product /\nTechnology / Event", fill=WHITE, size=13)
    _rounded_box(slide, 8.28, 2.65, 1.98, 0.85, "常识实体层\nCompany / Organization /\nPerson / ProductObject", fill=WHITE, size=13)
    _rounded_box(slide, 8.28, 3.75, 1.98, 0.8, "事件层\nEvent / Policy /\nCooperation / Financing", fill=WHITE, size=13)
    _rounded_box(slide, 8.28, 4.8, 1.98, 0.8, "证据层\nDocument /\nChunk / DataSource", fill=WHITE, size=13)

    # Outputs
    _rounded_box(slide, 10.95, 1.75, 1.75, 0.95, "产业网链推理\n上下游 / 区域 /\n技术 / 主体传播", fill=WHITE, size=13)
    _rounded_box(slide, 10.95, 3.0, 1.75, 0.95, "事件传导分析\n影响范围 / 关联路径 /\n风险扩散", fill=WHITE, size=13)
    _rounded_box(slide, 10.95, 4.25, 1.75, 0.95, "抽取增强与应用\nKAG / OpenKS 约束\n问答 / 预警 / 溯源", fill=WHITE, size=13)

    # Arrows between columns
    _arrow(slide, 2.32, 2.25, 2.93, 2.25)
    _arrow(slide, 2.32, 3.5, 2.93, 3.5)
    _arrow(slide, 2.32, 4.75, 2.93, 4.75)

    _arrow(slide, 4.62, 2.25, 5.28, 1.95)
    _arrow(slide, 4.62, 3.5, 5.28, 2.95)
    _arrow(slide, 4.62, 4.75, 5.28, 4.95)

    _arrow(slide, 7.62, 1.94, 8.25, 1.98)
    _arrow(slide, 7.62, 2.94, 8.25, 3.0)
    _arrow(slide, 7.62, 3.94, 8.25, 4.08)
    _arrow(slide, 7.62, 4.94, 8.25, 5.18)

    _arrow(slide, 10.28, 1.98, 10.92, 2.18)
    _arrow(slide, 10.28, 4.08, 10.92, 3.48)
    _arrow(slide, 10.28, 5.18, 10.92, 4.73)

    # Internal graph vertical hint
    _arrow(slide, 9.27, 2.4, 9.27, 2.63, color=TEAL, width=1.8)
    _arrow(slide, 9.27, 3.5, 9.27, 3.73, color=TEAL, width=1.8)
    _arrow(slide, 9.27, 4.55, 9.27, 4.78, color=TEAL, width=1.8)

    legend = slide.shapes.add_textbox(Inches(0.45), Inches(6.55), Inches(12.3), Inches(0.35))
    _set_text(
        legend,
        "核心逻辑：常识层提供背景骨架，事实层提供动态事件，融合层统一 schema 和主实体，统一大图同时服务推理、抽取增强和问答应用。",
        size=14,
        bold=True,
        color=NAVY,
        valign=MSO_ANCHOR.TOP,
    )


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover(prs)
    overview(prs)

    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(f"written {OUT_PPT}")


if __name__ == "__main__":
    build()
